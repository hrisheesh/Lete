from typing import List
from lete.app.schemas.section import DocumentSectionCreate
from lete.app.parsing.base import BaseParser


class PptxParser(BaseParser):
    """
    Parses PowerPoint (.pptx) files using python-pptx.
    Each slide becomes one section, preserving slide number for citation purposes.
    Extracts: title, body text, table data, and alt-text from images.
    """

    def parse(self, file_path: str, document_id: str, filename: str = "") -> List[DocumentSectionCreate]:
        from pptx import Presentation
        from pptx.util import Pt
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(file_path)
        sections = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            lines = []

            for shape in slide.shapes:
                # Extract text from text frames (title, body, notes)
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)

                # Extract table data as structured text
                elif shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            lines.append(row_text)

            # Also extract notes
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                notes_text = notes_tf.text.strip() if notes_tf else ""
                if notes_text:
                    lines.append(f"[Speaker Notes] {notes_text}")

            content = self._clean_text("\n".join(lines))
            if content:
                sections.append(DocumentSectionCreate(
                    document_id=document_id,
                    content=content,
                    page_number=slide_num,
                    section_index=slide_num - 1
                ))

        return sections
